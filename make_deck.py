"""Generate a PowerPoint deck with one slide per datacenter GPU chart."""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os

CURVES_DIR = "outputs/curves"
OUTPUT_PATH = "outputs/datacenter_gpu_deck.pptx"

# Load model results
df = pd.read_csv("outputs/curve_fit_results.csv")
dc = df[df["segment"] == "DATACENTER"].sort_values("proj_yr5_ratio")

# Load ranking for MSRP
rk = pd.read_csv("outputs/depreciation_ranking.csv")
msrp_map = dict(zip(rk["gpu_name"], rk["msrp_usd"]))

def add_image_preserved_ratio(slide, fpath, max_width_in, max_height_in, left_in, top_in):
    """Add image to slide preserving aspect ratio within max bounds, centered."""
    img = Image.open(fpath)
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Fit within max bounds
    w = max_width_in
    h = w / aspect
    if h > max_height_in:
        h = max_height_in
        w = h * aspect

    # Center horizontally
    actual_left = left_in + (max_width_in - w) / 2
    slide.shapes.add_picture(fpath, Inches(actual_left), Inches(top_in),
                             Inches(w), Inches(h))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Datacenter GPU Depreciation Curves"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = f"{len(dc)} GPUs | Ranked by 5-Year Terminal Value Ratio (TVR)"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "March 2026"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p3.alignment = PP_ALIGN.CENTER

# Summary methodology slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

p_title = tf.paragraphs[0]
p_title.text = "Methodology"
p_title.font.size = Pt(32)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p_title.alignment = PP_ALIGN.LEFT

p_body = tf.add_paragraph()
p_body.space_before = Pt(20)
p_body.text = (
    f"We modeled depreciation curves for {len(dc)} datacenter GPUs using "
    "over 134,000 primary and 8.7 million secondary price observations. "
    "Each GPU's price history was normalized to its launch MSRP and fit "
    "with exponential and power-law depreciation curves. The probability "
    "of falling below key value thresholds was estimated via 10,000 "
    "Monte Carlo simulations, producing the confidence intervals shown "
    "on each chart. All projections extend to ten years from launch."
)
p_body.font.size = Pt(16)
p_body.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p_body.line_spacing = Pt(24)

for _, row in dc.iterrows():
    gpu = row["gpu_name"]
    tvr = row["proj_yr5_ratio"]
    n_months = int(row["n_months"]) if pd.notna(row["n_months"]) else 0
    carrier = row.get("is_carrier_adjusted", False)
    msrp = msrp_map.get(gpu, 0)
    conf = row["conf_flag"]

    fname = gpu.replace(" ", "_").replace("(", "").replace(")", "") + ".png"
    fpath = os.path.join(CURVES_DIR, fname)
    if not os.path.exists(fpath):
        print(f"  SKIP (no chart): {gpu} -> {fpath}")
        continue

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add chart image — preserve aspect ratio
    add_image_preserved_ratio(slide, fpath,
                              max_width_in=11.7, max_height_in=5.0,
                              left_in=0.8, top_in=0.2)

    # Text box below chart
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Line 1: GPU name and key stats
    tvr_pct = f"{tvr*100:.1f}%"
    dollar_val = f"${tvr * msrp:,.0f}" if msrp > 0 else ""
    loss_pct = f"{(1-tvr)*100:.1f}%" if tvr <= 1.0 else f"+{(tvr-1)*100:.1f}%"

    p1 = tf.paragraphs[0]
    p1.text = f"{gpu}  |  TVR {tvr_pct} ({dollar_val})  |  {loss_pct} loss at 5yr"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    print(f"  Added: {gpu} (TVR {tvr_pct})")

# Family and segment charts
family_dir = os.path.join(CURVES_DIR, "families")
segment_dir = os.path.join(CURVES_DIR, "segments")

for chart_dir, label in [(family_dir, "Family"), (segment_dir, "Segment")]:
    if not os.path.isdir(chart_dir):
        continue
    for fname in sorted(os.listdir(chart_dir)):
        if not fname.endswith(".png"):
            continue
        dc_keywords = ["V100", "A100", "H100", "H200", "Datacenter"]
        if not any(kw in fname for kw in dc_keywords):
            continue

        fpath = os.path.join(chart_dir, fname)
        title = fname.replace("_", " ").replace(".png", "")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_image_preserved_ratio(slide, fpath,
                                  max_width_in=11.7, max_height_in=5.0,
                                  left_in=0.8, top_in=0.2)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.0))
        tf = txBox.text_frame
        p1 = tf.paragraphs[0]
        p1.text = f"{title}  |  {label} Overview"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        print(f"  Added {label}: {title}")

prs.save(OUTPUT_PATH)
print(f"\nSaved {OUTPUT_PATH} ({len(prs.slides)} slides)")
